"""
ぱむぱむ プレゼンテーション PPTX 生成スクリプト
Midnight Galaxy テーマ
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── カラーパレット (Midnight Galaxy) ──────────────────
C_DEEP_PURPLE  = RGBColor(0x2b, 0x1e, 0x3e)
C_DARK_BG      = RGBColor(0x1a, 0x10, 0x28)
C_DEMO_BG      = RGBColor(0x12, 0x0d, 0x1e)
C_COSMIC_BLUE  = RGBColor(0x4a, 0x4e, 0x8f)
C_LAVENDER     = RGBColor(0xa4, 0x90, 0xc2)
C_SILVER       = RGBColor(0xe6, 0xe6, 0xfa)
C_ACCENT_PINK  = RGBColor(0xe8, 0x79, 0xa0)
C_ACCENT_GOLD  = RGBColor(0xf4, 0xd0, 0x3f)
C_WHITE        = RGBColor(0xff, 0xff, 0xff)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank


def add_slide(bg_color=None):
    slide = prs.slides.add_slide(blank_layout)
    if bg_color:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return slide


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=24, bold=False, color=None,
             align=PP_ALIGN.LEFT, v_anchor=None, wrap=True):
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or C_SILVER
    run.font.name = "Helvetica Neue"
    return txBox


def add_multiline_text(slide, lines, left, top, width, height,
                       font_size=20, color=None, bold=False, align=PP_ALIGN.LEFT, spacing_pt=None):
    """lines = [(text, size, bold, color), ...]"""
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            text, size, b, c = item, font_size, bold, color or C_SILVER
        else:
            text, size, b, c = item

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = align
        if spacing_pt:
            p.space_before = Pt(spacing_pt)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = b
        run.font.color.rgb = c
        run.font.name = "Helvetica Neue"
    return txBox


def add_pill(slide, text, left, top, color=C_ACCENT_PINK, bg_alpha_fill=None):
    w = Inches(2.2)
    h = Inches(0.35)
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x4a, 0x4e, 0x8f) if color == C_COSMIC_BLUE else RGBColor(0x3a, 0x1a, 0x2c)
    shape.line.color.rgb = color
    shape.line.width = Pt(1)
    # Corner rounding
    from pptx.oxml.ns import qn
    sp_pr = shape._element.spPr
    prstGeom = sp_pr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            from lxml import etree
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        from lxml import etree
        gd = etree.SubElement(avLst, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', 'val 50000')

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "Helvetica Neue"
    return shape


def slide_number(slide, n, total=9):
    add_text(slide, f"{n} / {total}",
             SLIDE_W - Inches(1.2), Pt(20), Inches(1.0), Inches(0.3),
             font_size=9, color=C_LAVENDER, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# Slide 1 — Title
# ══════════════════════════════════════════════════════════════
s1 = add_slide(C_DARK_BG)
slide_number(s1, 1)

# Decorative orb (blurred circle)
orb = add_rect(s1, Inches(9.5), Inches(-1), Inches(5), Inches(5), fill_color=RGBColor(0x3a, 0x10, 0x28))
orb.line.fill.background()

add_pill(s1, "✦ ハッカソン作品発表", Inches(0.7), Inches(1.4), C_ACCENT_PINK)

add_multiline_text(s1, [
    ("ぱむぱむ", 64, True, C_SILVER),
    ("〜 AI 手相占い 〜", 26, False, C_LAVENDER),
], left=Inches(0.7), top=Inches(1.9), width=Inches(9), height=Inches(2.8))

add_multiline_text(s1, [
    ("振動センサーで動揺を検知し、", 18, False, C_SILVER),
    ("Gemini AI が容赦なく追い込む体験型アプリ", 18, False, C_SILVER),
], left=Inches(0.7), top=Inches(4.3), width=Inches(9), height=Inches(1.0), spacing_pt=4)

# Tech pills row
for i, label in enumerate(["Raspberry Pi", "Gemini Live API", "FastAPI + React"]):
    x = Inches(0.7) + i * Inches(2.5)
    b = add_rect(s1, x, Inches(5.6), Inches(2.2), Inches(0.45),
                 fill_color=RGBColor(0x2a, 0x2a, 0x50), line_color=C_COSMIC_BLUE)
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(12)
    r.font.color.rgb = C_LAVENDER
    r.font.name = "Helvetica Neue"


# ══════════════════════════════════════════════════════════════
# Slide 2 — なぜ占いは当たるのか
# ══════════════════════════════════════════════════════════════
s2 = add_slide(C_DEEP_PURPLE)
slide_number(s2, 2)
add_pill(s2, "✦ きっかけ", Inches(0.7), Inches(0.5), C_ACCENT_GOLD)

add_multiline_text(s2, [
    ('なぜ手相占いは「当たる」のか？', 36, True, C_SILVER),
], left=Inches(0.7), top=Inches(0.9), width=Inches(11), height=Inches(0.9))

# Table header bar
add_rect(s2, Inches(0.7), Inches(2.0), Inches(11.8), Inches(0.4),
         fill_color=RGBColor(0x3a, 0x2a, 0x52))

cols = [Inches(0.7), Inches(5.5), Inches(6.5)]
headers = ["本物の占い師がやっていること", "→", "ぱむぱむでの実現"]
widths  = [Inches(4.5), Inches(0.6), Inches(5.4)]
for j, (hd, cx, cw) in enumerate(zip(headers, cols, widths)):
    add_text(s2, hd, cx, Inches(2.05), cw, Inches(0.35),
             font_size=10, bold=True,
             color=C_LAVENDER if j != 1 else C_ACCENT_PINK,
             align=PP_ALIGN.CENTER if j == 1 else PP_ALIGN.LEFT)

rows = [
    ("広い仮説を投げ、反応を観察する（バーナム効果）",
     "Stage1：誰にでも当てはまる感情仮説を投下"),
    ("表情・呼吸・手の震えで反応を読む（Cold Reading）",
     "振動センサーで手の震え = 動揺率（0〜100%）に変換"),
    ("反応が出たら「図星」として断言・追い込む",
     "Stage2：動揺率に応じてAIの「圧」を最大化"),
]
for i, (left_txt, right_txt) in enumerate(rows):
    y = Inches(2.45) + i * Inches(0.85)
    bg = RGBColor(0x28, 0x1c, 0x3c) if i % 2 == 0 else RGBColor(0x24, 0x18, 0x36)
    add_rect(s2, Inches(0.7), y, Inches(11.8), Inches(0.8), fill_color=bg)
    add_text(s2, left_txt,  Inches(0.8), y + Pt(8), Inches(4.4), Inches(0.7), font_size=11, color=C_LAVENDER)
    add_text(s2, "→",        Inches(5.4), y + Pt(8), Inches(0.6), Inches(0.7), font_size=14, color=C_ACCENT_PINK, align=PP_ALIGN.CENTER)
    add_text(s2, right_txt, Inches(6.1), y + Pt(8), Inches(6.2), Inches(0.7), font_size=11, color=C_SILVER)

# Quote
add_rect(s2, Inches(0.7), Inches(5.45), Inches(0.06), Inches(0.7), fill_color=C_ACCENT_PINK)
add_text(s2, '「瞳孔の反応を読む」をセンサーで代替できるんじゃないか',
         Inches(0.9), Inches(5.45), Inches(11), Inches(0.75),
         font_size=16, bold=True, color=C_SILVER)


# ══════════════════════════════════════════════════════════════
# Slide 3 — コンセプト
# ══════════════════════════════════════════════════════════════
s3 = add_slide(C_DARK_BG)
slide_number(s3, 3)
add_pill(s3, "✦ コンセプト", Inches(0.7), Inches(0.5), C_ACCENT_PINK)

add_multiline_text(s3, [
    ("振動センサーで動揺を検知し、", 32, True, C_SILVER),
    ("AI が追い込む", 32, True, C_ACCENT_PINK),
], left=Inches(0.7), top=Inches(0.9), width=Inches(11), height=Inches(1.4))

# Flow boxes
flow_items = [
    ("🖐️", "手を乗せる"),
    ("→", ""),
    ("📡", "振動センサー\n（ラズパイ）"),
    ("→", ""),
    ("📊", "動揺率\n0〜100%"),
    ("→", ""),
    ("🤖", "Gemini が\n追い込む"),
]
box_w = Inches(1.5)
box_h = Inches(1.3)
start_x = Inches(0.6)
y = Inches(2.5)
arrow_w = Inches(0.6)

x = start_x
for icon, label in flow_items:
    if icon == "→":
        add_text(s3, "→", x, y + Inches(0.4), arrow_w, Inches(0.5),
                 font_size=22, bold=True, color=C_ACCENT_PINK, align=PP_ALIGN.CENTER)
        x += arrow_w
    else:
        is_last = (icon == "🤖")
        fill = RGBColor(0x3a, 0x18, 0x2c) if is_last else RGBColor(0x2a, 0x2a, 0x50)
        line = C_ACCENT_PINK if is_last else C_COSMIC_BLUE
        b = add_rect(s3, x, y, box_w, box_h, fill_color=fill, line_color=line, line_width=Pt(1.5))
        add_multiline_text(s3, [
            (icon, 26, False, C_SILVER),
            (label, 11, False, C_LAVENDER if not is_last else C_ACCENT_PINK),
        ], left=x, top=y, width=box_w, height=box_h, align=PP_ALIGN.CENTER)
        x += box_w

# Two stat boxes
stat_y = Inches(4.2)
for i, (val, lbl, sub, col) in enumerate([
    ("0 → 100%", "動揺率", "スライディングウィンドウ集計", C_ACCENT_GOLD),
    ("探り → 断言", "AI の態度", "Cold Reading のロジックを実装", C_ACCENT_PINK),
]):
    bx = Inches(0.7) + i * Inches(5.7)
    add_rect(s3, bx, stat_y, Inches(5.3), Inches(1.3),
             fill_color=RGBColor(0x24, 0x24, 0x48), line_color=C_COSMIC_BLUE)
    add_text(s3, lbl, bx + Inches(0.15), stat_y + Inches(0.1), Inches(5), Inches(0.3),
             font_size=10, color=C_LAVENDER)
    add_text(s3, val, bx + Inches(0.15), stat_y + Inches(0.32), Inches(5), Inches(0.55),
             font_size=22, bold=True, color=col)
    add_text(s3, sub, bx + Inches(0.15), stat_y + Inches(0.9), Inches(5), Inches(0.3),
             font_size=9, color=C_LAVENDER)


# ══════════════════════════════════════════════════════════════
# Slide 4 — Demo 1
# ══════════════════════════════════════════════════════════════
s4 = add_slide(C_DEMO_BG)
slide_number(s4, 4)

add_text(s4, "🖐️", Inches(5.9), Inches(0.8), Inches(1.5), Inches(1.2),
         font_size=56, align=PP_ALIGN.CENTER)
add_text(s4, "Demo 1", Inches(3), Inches(2.0), Inches(7.3), Inches(0.8),
         font_size=40, bold=True, color=C_ACCENT_PINK, align=PP_ALIGN.CENTER)
add_text(s4, "実際に占われてみる", Inches(2), Inches(2.85), Inches(9.3), Inches(0.65),
         font_size=28, bold=True, color=C_SILVER, align=PP_ALIGN.CENTER)
add_text(s4, "手を乗せると… AIが喋り始める → 体が反応する → AIが追い込む",
         Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.55),
         font_size=16, color=C_LAVENDER, align=PP_ALIGN.CENTER)

for i, (label, col) in enumerate([("約 2 分", C_COSMIC_BLUE), ("実機デモ動画", C_ACCENT_PINK)]):
    bx = Inches(4.9) + i * Inches(2.0)
    b = add_rect(s4, bx, Inches(4.5), Inches(1.7), Inches(0.42),
                 fill_color=RGBColor(0x1a, 0x1a, 0x3a), line_color=col)
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = col
    r.font.name = "Helvetica Neue"


# ══════════════════════════════════════════════════════════════
# Slide 5 — アーキテクチャ
# ══════════════════════════════════════════════════════════════
s5 = add_slide(C_DEEP_PURPLE)
slide_number(s5, 5)
add_pill(s5, "✦ 仕組み", Inches(0.7), Inches(0.5), C_COSMIC_BLUE)
add_text(s5, "どうやって動いているか", Inches(0.7), Inches(0.92), Inches(11), Inches(0.75),
         font_size=34, bold=True, color=C_SILVER)

arch = [
    ("📡", "Raspberry Pi", "振動センサー\n0/1 パルス", False),
    ("⚙️", "バックエンド", "FastAPI\n動揺率エンジン", False),
    ("🤖", "Gemini Live", "リアルタイム\n音声生成", True),
    ("🌐", "ブラウザ", "React\n音声再生", False),
]
arrow_labels = ["WebSocket", "Tool Use", "TTS 音声"]

bw = Inches(2.0)
bh = Inches(1.5)
aw = Inches(1.0)
total_w = len(arch) * bw + len(arrow_labels) * aw
start_x = (SLIDE_W - total_w) / 2
ay = Inches(2.1)

x = start_x
for idx, (icon, name, detail, is_main) in enumerate(arch):
    fill = RGBColor(0x3a, 0x18, 0x2c) if is_main else RGBColor(0x22, 0x18, 0x38)
    line = C_ACCENT_PINK if is_main else C_COSMIC_BLUE
    add_rect(s5, x, ay, bw, bh, fill_color=fill, line_color=line, line_width=Pt(1.5))
    add_multiline_text(s5, [
        (icon, 24, False, C_SILVER),
        (name, 13, True, C_SILVER),
        (detail, 10, False, C_LAVENDER),
    ], left=x, top=ay, width=bw, height=bh, align=PP_ALIGN.CENTER)
    x += bw
    if idx < len(arrow_labels):
        add_text(s5, "→", x, ay + Inches(0.5), aw * 0.6, Inches(0.4),
                 font_size=20, color=C_ACCENT_PINK, align=PP_ALIGN.CENTER)
        add_text(s5, arrow_labels[idx], x, ay + Inches(0.95), aw, Inches(0.3),
                 font_size=8, color=C_LAVENDER, align=PP_ALIGN.CENTER)
        x += aw

# Two info boxes
info_y = Inches(3.9)
for i, (title, body) in enumerate([
    ("動揺率エンジン",
     "直近 10 秒の振動数をスライディングウィンドウで集計。急上昇（+30%）検知時にAIへプッシュ通知。"),
    ("Gemini Tool Use",
     "AI が自発的に get_agitation() を呼び出し、リアルタイムで動揺率を取得しながら応答を組み立てる。"),
]):
    bx = Inches(0.7) + i * Inches(6.1)
    add_rect(s5, bx, info_y, Inches(5.7), Inches(1.3),
             fill_color=RGBColor(0x24, 0x24, 0x44), line_color=C_COSMIC_BLUE)
    add_text(s5, title, bx + Inches(0.15), info_y + Inches(0.1), Inches(5.4), Inches(0.3),
             font_size=10, bold=True, color=C_LAVENDER)
    add_text(s5, body, bx + Inches(0.15), info_y + Inches(0.38), Inches(5.4), Inches(0.85),
             font_size=11, color=C_SILVER)


# ══════════════════════════════════════════════════════════════
# Slide 6 — AIの追い込みロジック
# ══════════════════════════════════════════════════════════════
s6 = add_slide(C_DARK_BG)
slide_number(s6, 6)
add_pill(s6, "✦ AI の追い込みロジック", Inches(0.7), Inches(0.5), C_ACCENT_PINK)
add_multiline_text(s6, [
    ("動揺率で変わる、", 32, True, C_SILVER),
    ("AI の圧", 32, True, C_ACCENT_PINK),
], left=Inches(0.7), top=Inches(0.9), width=Inches(11), height=Inches(1.0))

headers6 = ["動揺率", "強度", "AI の台詞（例）"]
col_x6 = [Inches(0.7), Inches(2.8), Inches(4.5)]
col_w6 = [Inches(1.9), Inches(1.5), Inches(8.1)]

# Header row
add_rect(s6, Inches(0.7), Inches(2.05), Inches(12.3), Inches(0.38),
         fill_color=RGBColor(0x2a, 0x20, 0x40))
for j, (hd, cx, cw) in enumerate(zip(headers6, col_x6, col_w6)):
    add_text(s6, hd, cx + Inches(0.1), Inches(2.1), cw, Inches(0.3),
             font_size=9, bold=True, color=C_LAVENDER, align=PP_ALIGN.LEFT)

rows6 = [
    ("0 〜 30%",  C_COSMIC_BLUE,  0.28, "探り・暗示",   C_LAVENDER,  "「何かが微かに動いています……何を思い出しましたか？」",  C_SILVER),
    ("30 〜 60%", C_LAVENDER,     0.56, "確信・絞り込み", C_LAVENDER, "「体が反応しました。その感覚、ずっと持っていたのでは？」", C_SILVER),
    ("60 〜 80%", C_ACCENT_PINK,  0.84, "断言・追い込み", C_ACCENT_PINK, "「それは○○への恐れです——言えなかった言葉があるでしょう？」", C_ACCENT_PINK),
    ("80% 〜",    C_ACCENT_GOLD,  1.0,  "完全断言・畳み掛け", C_ACCENT_GOLD, "「隠せていません。その名前、今頭に浮かんでいますね？」", C_ACCENT_GOLD),
]

for i, (level_txt, bar_color, bar_ratio, intensity, int_col, speech, speech_col) in enumerate(rows6):
    ry = Inches(2.45) + i * Inches(0.9)
    bg = RGBColor(0x1e, 0x16, 0x30) if i % 2 == 0 else RGBColor(0x22, 0x1a, 0x36)
    add_rect(s6, Inches(0.7), ry, Inches(12.3), Inches(0.85), fill_color=bg)

    # level text
    add_text(s6, level_txt, col_x6[0] + Inches(0.1), ry + Inches(0.05),
             Inches(1.3), Inches(0.4), font_size=11, color=C_SILVER)
    # bar
    bar_total_w = Inches(1.5)
    add_rect(s6, col_x6[0] + Inches(0.1), ry + Inches(0.5),
             Inches(bar_total_w.inches * bar_ratio), Inches(0.12), fill_color=bar_color)

    # intensity
    add_text(s6, intensity, col_x6[1] + Inches(0.05), ry + Inches(0.2),
             col_w6[1], Inches(0.5), font_size=11, color=int_col)

    # speech
    add_text(s6, speech, col_x6[2] + Inches(0.05), ry + Inches(0.15),
             col_w6[2] - Inches(0.2), Inches(0.55), font_size=12, color=speech_col)


# ══════════════════════════════════════════════════════════════
# Slide 7 — Demo 2
# ══════════════════════════════════════════════════════════════
s7 = add_slide(C_DEMO_BG)
slide_number(s7, 7)

add_text(s7, "🔥", Inches(5.9), Inches(0.8), Inches(1.5), Inches(1.2),
         font_size=56, align=PP_ALIGN.CENTER)
add_text(s7, "Demo 2", Inches(3), Inches(2.0), Inches(7.3), Inches(0.8),
         font_size=40, bold=True, color=C_ACCENT_PINK, align=PP_ALIGN.CENTER)
add_text(s7, "動揺 MAX → AI が追い込むシーン", Inches(2), Inches(2.85), Inches(9.3), Inches(0.65),
         font_size=28, bold=True, color=C_SILVER, align=PP_ALIGN.CENTER)
add_text(s7, "動揺率が上昇するにつれてAIの「圧」が変わる様子",
         Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.55),
         font_size=16, color=C_LAVENDER, align=PP_ALIGN.CENTER)

for i, (label, col) in enumerate([("約 2 分", C_COSMIC_BLUE), ("「隠せていません」", C_ACCENT_GOLD)]):
    bx = Inches(4.2) + i * Inches(2.6)
    b = add_rect(s7, bx, Inches(4.5), Inches(2.3), Inches(0.42),
                 fill_color=RGBColor(0x1a, 0x1a, 0x3a), line_color=col)
    tf = b.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = col
    r.font.name = "Helvetica Neue"


# ══════════════════════════════════════════════════════════════
# Slide 8 — 技術スタック
# ══════════════════════════════════════════════════════════════
s8 = add_slide(C_DEEP_PURPLE)
slide_number(s8, 8)
add_pill(s8, "✦ 技術スタック", Inches(0.7), Inches(0.5), C_COSMIC_BLUE)
add_text(s8, "使った技術", Inches(0.7), Inches(0.92), Inches(11), Inches(0.75),
         font_size=34, bold=True, color=C_SILVER)

tech_cards = [
    ("📡", "Raspberry Pi", "振動センサー\nGPIO 0/1 パルス", C_COSMIC_BLUE),
    ("🤖", "Gemini Live API", "gemini-live-2.5-flash\nTool Use + VAD", C_ACCENT_PINK),
    ("🗣️", "Gemini TTS", "gemini-2.5-flash-tts\nリアルタイム音声生成", C_ACCENT_PINK),
    ("⚙️", "FastAPI + React", "WebSocket 接続\nDocker / OrbStack", C_COSMIC_BLUE),
]
card_w = Inches(2.8)
card_h = Inches(1.8)
gap = Inches(0.35)
total_cards_w = 4 * card_w + 3 * gap
sx = (SLIDE_W - total_cards_w) / 2
cy = Inches(2.0)

for i, (icon, name, detail, col) in enumerate(tech_cards):
    cx = sx + i * (card_w + gap)
    fill = RGBColor(0x3a, 0x18, 0x2c) if col == C_ACCENT_PINK else RGBColor(0x22, 0x1e, 0x3a)
    add_rect(s8, cx, cy, card_w, card_h, fill_color=fill, line_color=col, line_width=Pt(1.5))
    add_multiline_text(s8, [
        (icon, 28, False, C_SILVER),
        (name, 14, True, C_SILVER),
        (detail, 10, False, C_LAVENDER),
    ], left=cx, top=cy, width=card_w, height=card_h, align=PP_ALIGN.CENTER)

# Info box
add_rect(s8, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.3),
         fill_color=RGBColor(0x22, 0x22, 0x3e), line_color=C_COSMIC_BLUE)
add_text(s8, "ポイント：Gemini が自分で動揺率を取りに行く",
         Inches(0.9), Inches(4.18), Inches(11.6), Inches(0.3),
         font_size=10, bold=True, color=C_LAVENDER)
add_text(s8,
         "AI はバックエンドの get_agitation() ツールを自律的に呼び出し、リアルタイムの動揺率を取得しながら応答を生成する。センサーと言語モデルがリアルタイムに連携。",
         Inches(0.9), Inches(4.5), Inches(11.6), Inches(0.85),
         font_size=13, color=C_SILVER)


# ══════════════════════════════════════════════════════════════
# Slide 9 — まとめ
# ══════════════════════════════════════════════════════════════
s9 = add_slide(C_DARK_BG)
slide_number(s9, 9)
add_pill(s9, "✦ まとめ", Inches(5.6), Inches(0.5), C_ACCENT_PINK)

add_multiline_texts = [
    ("占いはフィクションかもしれない。", 28, False, C_SILVER),
    ("でも、あなたの", 28, False, C_SILVER),
    ("体の反応", 28, True, C_ACCENT_PINK),
    ("は", 28, False, C_SILVER),
]
# single line with mix: do two text boxes
add_text(s9, "占いはフィクションかもしれない。", Inches(2), Inches(1.6), Inches(9.3), Inches(0.7),
         font_size=28, bold=False, color=C_SILVER, align=PP_ALIGN.CENTER)
add_multiline_text(s9, [
    ("でも、あなたの体の反応は", 28, False, C_SILVER),
    ("本物だ。", 36, True, C_ACCENT_GOLD),
], left=Inches(2), top=Inches(2.3), width=Inches(9.3), height=Inches(1.6),
   align=PP_ALIGN.CENTER, spacing_pt=6)

# Three stats
stats = [
    ("2", "API の活用\n（Live + TTS）", C_ACCENT_PINK),
    ("4", "ハードウェア\n+ ソフト連携", C_ACCENT_GOLD),
    ("∞", "AI に\n追い込まれる体験", C_SILVER),
]
stat_w = Inches(2.8)
stat_sx = (SLIDE_W - 3 * stat_w - 2 * Inches(0.5)) / 2
for i, (num, lbl, col) in enumerate(stats):
    sx2 = stat_sx + i * (stat_w + Inches(0.5))
    add_text(s9, num, sx2, Inches(4.2), stat_w, Inches(0.7),
             font_size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s9, lbl, sx2, Inches(4.9), stat_w, Inches(0.6),
             font_size=11, color=C_LAVENDER, align=PP_ALIGN.CENTER)

# Dividers
for i in range(2):
    dx = stat_sx + (i + 1) * (stat_w + Inches(0.5)) - Inches(0.3)
    add_rect(s9, dx, Inches(4.2), Inches(0.02), Inches(1.2),
             fill_color=C_COSMIC_BLUE)

add_text(s9, "ご清聴ありがとうございました",
         Inches(2), Inches(6.2), Inches(9.3), Inches(0.6),
         font_size=20, color=C_LAVENDER, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────
out = "/Users/kubotadaichi/dev/github/PalmPalm/docs/slides/presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
